#!/usr/bin/env python3
"""
Convert lecture materials (PDF, DOCX, PPTX) to JSON flashcards and quiz questions.
Usage: python convert_materials.py --input ../Material/KLR --output ../data/klr-fibu --subject klr
"""
import argparse
import json
import os
import re
import sys
from pathlib import Path

def extract_text_pdf(path):
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)
    except ImportError:
        print("Install pdfplumber: pip install pdfplumber", file=sys.stderr)
        return ""

def extract_text_docx(path):
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except ImportError:
        print("Install python-docx: pip install python-docx", file=sys.stderr)
        return ""

def extract_text_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)
    except ImportError:
        print("Install python-pptx: pip install python-pptx", file=sys.stderr)
        return ""

def extract_text(path):
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return extract_text_pdf(path)
    elif ext == ".docx":
        return extract_text_docx(path)
    elif ext == ".pptx":
        return extract_text_pptx(path)
    elif ext == ".rtf":
        with open(path, "rb") as f:
            raw = f.read().decode("latin-1", errors="ignore")
        return re.sub(r"\\[a-z]+\d* ?|[{}]", "", raw)
    return ""

def make_flashcards(lines, subject):
    cards = []
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if len(line) > 10 and line.endswith("?"):
            front = line
            back_lines = []
            j = i + 1
            while j < len(lines) and len(back_lines) < 3:
                candidate = lines[j].strip()
                if candidate:
                    back_lines.append(candidate)
                j += 1
            if back_lines:
                idx = len(cards) + 1
                cards.append({
                    "id": f"{subject}-gen{idx:03d}",
                    "front": front,
                    "back": " ".join(back_lines),
                    "tags": ["Generiert"],
                    "latex": False
                })
                i = j
                continue
        i += 1
    return cards

def make_quiz(lines, subject):
    questions = []
    for i, line in enumerate(lines):
        line = line.strip()
        if len(line) > 15 and line.endswith("?"):
            opts = []
            j = i + 1
            while j < len(lines) and len(opts) < 4:
                candidate = lines[j].strip()
                if 3 < len(candidate) < 120 and not candidate.endswith("?"):
                    opts.append(candidate)
                j += 1
            if len(opts) >= 2:
                idx = len(questions) + 1
                questions.append({
                    "id": f"{subject}-qgen{idx:03d}",
                    "question": line,
                    "options": opts[:4],
                    "correct": 0,
                    "explanation": "Bitte korrekten Index und Erklärung manuell eintragen."
                })
    return questions

def process_dir(input_dir, output_dir, subject):
    input_path = Path(input_dir)
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    all_text = []
    for f in sorted(input_path.rglob("*")):
        if f.suffix.lower() in (".pdf", ".docx", ".pptx", ".rtf"):
            print(f"  Reading: {f.name}")
            text = extract_text(f)
            if text:
                all_text.append(text)

    if not all_text:
        print("No readable files found.")
        return

    lines = "\n".join(all_text).splitlines()
    lines = [l for l in lines if l.strip()]

    flashcards = make_flashcards(lines, subject)
    quiz = make_quiz(lines, subject)

    fc_path = output_path / "flashcards_generated.json"
    qz_path = output_path / "quiz_generated.json"

    with open(fc_path, "w", encoding="utf-8") as f:
        json.dump(flashcards, f, ensure_ascii=False, indent=2)
    with open(qz_path, "w", encoding="utf-8") as f:
        json.dump(quiz, f, ensure_ascii=False, indent=2)

    print(f"\n✓ {len(flashcards)} Karteikarten → {fc_path}")
    print(f"✓ {len(quiz)} Quizfragen → {qz_path}")
    print("\nHinweis: quiz_generated.json enthält Platzhalter für 'correct'-Index und 'explanation'. Bitte manuell prüfen.")

def main():
    parser = argparse.ArgumentParser(description="Convert lecture materials to JSON")
    parser.add_argument("--input", required=True, help="Input directory with PDF/DOCX/PPTX files")
    parser.add_argument("--output", required=True, help="Output directory for JSON files")
    parser.add_argument("--subject", required=True, help="Subject prefix for IDs (e.g. klr, it, mathe)")
    args = parser.parse_args()
    process_dir(args.input, args.output, args.subject)

if __name__ == "__main__":
    main()
